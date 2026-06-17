"""
Reproducible generator for docs/presentations/mpl_library.pptx.

Run from the repository root:
    python docs/presentations/generate_mpl_presentation.py

Requires: python-pptx >= 1.0
"""

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from pptx.util import Cm

# ---------------------------------------------------------------------------
# Output path
# ---------------------------------------------------------------------------
OUTPUT = Path(__file__).parent / "mpl_library.pptx"

# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_BG      = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK    = RGBColor(0x1A, 0x1A, 0x1A)
C_BLUE    = RGBColor(0x1B, 0x4F, 0x8A)   # primary accent
C_TEAL    = RGBColor(0x1A, 0x6E, 0x7D)   # secondary accent
C_GRAY    = RGBColor(0x55, 0x55, 0x55)   # body text
C_LGRAY   = RGBColor(0xDD, 0xDD, 0xDD)   # light rule / divider
C_MGRAY   = RGBColor(0x99, 0x99, 0x99)   # footer
C_BOX_BG  = RGBColor(0xEA, 0xF0, 0xF9)   # light blue fill for boxes
C_BOX_B   = RGBColor(0xA0, 0xBA, 0xD9)   # box border
C_TEAL_BG = RGBColor(0xE0, 0xF4, 0xF6)
C_TEAL_B  = RGBColor(0x7B, 0xBE, 0xC6)
C_GREEN   = RGBColor(0x1A, 0x6B, 0x3A)   # "future" items
C_AMBER   = RGBColor(0x7A, 0x52, 0x00)   # caution / limitation

# ---------------------------------------------------------------------------
# Slide dimensions: widescreen 13.33 × 7.5 in
# ---------------------------------------------------------------------------
SW = Inches(13.33)
SH = Inches(7.50)

FOOTER_TEXT = "MPL2030 — MPL loop simulation library"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def blank_slide(prs: Presentation):
    """Add a blank slide (layout 6) with white background."""
    layout = prs.slide_layouts[6]   # blank
    slide  = prs.slides.add_slide(layout)
    # Force white background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_BG
    return slide


def add_textbox(
    slide,
    left, top, width, height,
    text,
    size=18,
    bold=False,
    italic=False,
    color=None,
    align=PP_ALIGN.LEFT,
    wrap=True,
    font_name="Calibri",
):
    color = color or C_DARK
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = font_name
    return txBox


def add_multiline_textbox(
    slide,
    left, top, width, height,
    lines,            # list of (text, size, bold, italic, color, align)
    font_name="Calibri",
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (text, size, bold, italic, color, align) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size    = Pt(size)
        run.font.bold    = bold
        run.font.italic  = italic
        run.font.color.rgb = color
        run.font.name    = font_name
    return txBox


def add_rect(
    slide,
    left, top, width, height,
    fill_color=None,
    line_color=None,
    line_width=Pt(0.75),
):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height,
    )
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()

    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width     = line_width
    else:
        line.fill.background()
    return shape


def add_rect_with_text(
    slide,
    left, top, width, height,
    text,
    size=13,
    bold=False,
    fill_color=None,
    line_color=None,
    text_color=None,
    align=PP_ALIGN.CENTER,
    font_name="Calibri",
):
    text_color = text_color or C_DARK
    shape = add_rect(slide, left, top, width, height,
                     fill_color=fill_color, line_color=line_color)
    tf = shape.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = text_color
    run.font.name  = font_name
    return shape


def add_arrow_right(slide, left, top, width, height=Inches(0.08), color=None):
    """Add a simple right-pointing arrow shape."""
    color = color or C_GRAY
    from pptx.util import Pt as _Pt
    # Use a line connector
    cx = slide.shapes.add_connector(1, left, top, left + width, top)
    cx.line.color.rgb = color
    cx.line.width     = Pt(1.5)
    return cx


def add_arrow_down(slide, left, top, height, color=None):
    color = color or C_GRAY
    cx = slide.shapes.add_connector(1, left, top, left, top + height)
    cx.line.color.rgb = color
    cx.line.width     = Pt(1.5)
    return cx


def add_footer(slide, slide_number: int, total: int = 8):
    """Add footer bar at the bottom of the slide."""
    foot_h = Inches(0.32)
    foot_t = SH - foot_h

    # Light rule
    rule = add_rect(slide, 0, foot_t, SW, Inches(0.01),
                    fill_color=C_LGRAY, line_color=None)

    footer_str = f"{FOOTER_TEXT}    |    {slide_number} / {total}"
    add_textbox(
        slide,
        Inches(0.3), foot_t + Inches(0.04),
        SW - Inches(0.6), foot_h,
        footer_str,
        size=9,
        color=C_MGRAY,
        align=PP_ALIGN.CENTER,
    )


def set_notes(slide, text: str):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def title_bar(slide, title: str, subtitle: str = ""):
    """Standard title bar at the top of every content slide."""
    # Blue accent strip
    add_rect(slide, 0, 0, SW, Inches(0.08), fill_color=C_BLUE)
    # Title text
    add_textbox(
        slide,
        Inches(0.4), Inches(0.12),
        SW - Inches(0.8), Inches(0.65),
        title,
        size=26,
        bold=True,
        color=C_BLUE,
    )
    if subtitle:
        add_textbox(
            slide,
            Inches(0.4), Inches(0.72),
            SW - Inches(0.8), Inches(0.35),
            subtitle,
            size=14,
            italic=True,
            color=C_GRAY,
        )


def content_top(has_subtitle=False) -> float:
    """Y coordinate where content starts, in Inches."""
    return 1.10 if has_subtitle else 0.90


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------

def slide_01(prs):
    slide = blank_slide(prs)

    # Blue top bar
    add_rect(slide, 0, 0, SW, Inches(0.12), fill_color=C_BLUE)

    # Blue bottom decorative band
    add_rect(slide, 0, SH - Inches(1.2), SW, Inches(1.2), fill_color=C_BLUE)

    # Main title
    add_textbox(
        slide,
        Inches(0.6), Inches(0.5),
        SW - Inches(1.2), Inches(1.4),
        "A Modular Python Framework\nfor Mechanically Pumped Loop Simulation",
        size=34, bold=True, color=C_BLUE, align=PP_ALIGN.LEFT,
    )

    # Rule
    add_rect(slide, Inches(0.6), Inches(2.05), Inches(9.0), Inches(0.04),
             fill_color=C_BLUE)

    # Subtitle / motivation statement
    add_textbox(
        slide,
        Inches(0.6), Inches(2.2),
        Inches(10.5), Inches(1.4),
        (
            "MPL modelling requires physically transparent coupling of\n"
            "pressure, enthalpy, mass flow, phase change,\n"
            "accumulator pressure reference, and empirical closures."
        ),
        size=18, color=C_DARK, align=PP_ALIGN.LEFT,
    )

    # Context bullets
    add_multiline_textbox(
        slide,
        Inches(0.6), Inches(3.7),
        Inches(10.5), Inches(2.2),
        [
            ("Context and motivation", 14, True, False, C_BLUE, PP_ALIGN.LEFT),
            ("", 8, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  •  MPL2030 programme: thermal control of spacecraft and high-power electronics", 14, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  •  Two-phase loops couple compressible accumulator dynamics, pump hydraulics,", 14, False, False, C_DARK, PP_ALIGN.LEFT),
            ("     phase change in evaporator and condenser, and semi-empirical closure relations", 14, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  •  Physical transparency and replaceability of closure models are primary design goals", 14, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  •  Steady-state foundation first; dynamic seams declared and prepared", 14, False, False, C_DARK, PP_ALIGN.LEFT),
        ],
    )

    # Bottom band text
    add_textbox(
        slide,
        Inches(0.6), SH - Inches(1.05),
        SW - Inches(1.2), Inches(0.45),
        "Scientific committee presentation  ·  2026",
        size=13, color=C_WHITE, align=PP_ALIGN.LEFT,
    )
    add_textbox(
        slide,
        Inches(0.6), SH - Inches(0.65),
        SW - Inches(1.2), Inches(0.45),
        "MPL2030 — MPL loop simulation library    |    1 / 8",
        size=11, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.LEFT,
    )

    set_notes(slide, (
        "Welcome. This presentation describes the physical modelling framework "
        "we are developing for Mechanically Pumped Loop simulation within MPL2030.\n\n"
        "Our central challenge is that MPL loops tightly couple global pressure, "
        "local enthalpy transport, two-phase thermodynamics, a pressurising accumulator, "
        "and empirical closure relations — all of which interact. "
        "A monolithic implementation makes any one of those elements hard to replace or validate in isolation.\n\n"
        "The framework is designed to keep every physical assumption explicit and replaceable."
    ))
    add_footer(slide, 1)
    return slide


# ---------------------------------------------------------------------------
# Slide 2 — Physical modelling problem
# ---------------------------------------------------------------------------

def slide_02(prs):
    slide = blank_slide(prs)
    title_bar(slide, "The Physical Modelling Problem",
              "A single-loop MPL with parallel evaporator branches")

    # We draw a simplified MPL schematic using rectangles and connectors.
    # Layout: horizontal loop, left-to-right on two rows.

    Y1 = Inches(1.55)   # top row y-centre  (we draw boxes centred here)
    Y2 = Inches(4.20)   # bottom row y-centre
    BW = Inches(1.50)   # box width
    BH = Inches(0.55)   # box height
    bh2 = BH / 2

    # Helper: box centres (left col → right col)
    # Top row (left to right): Accumulator, Pump, Pipe, [split], Evaporator A
    # Bottom row: return Pipe, Condenser ← Pipe ← [merge], Evaporator B
    # Variables labelled between elements.

    def bx(cx, cy): return Inches(cx) - BW / 2
    def by(cy):     return Inches(cy) - BH / 2

    boxes_top = [
        (1.45, 1.55, "Accumulator\n(PCA/HCA)", C_TEAL_BG, C_TEAL),
        (3.40, 1.55, "Pump",                   C_BOX_BG,  C_BOX_B),
        (5.35, 1.55, "Pipe\n(supply)",          C_BOX_BG,  C_BOX_B),
        (8.80, 1.55, "Evaporator A\nQ_in",      C_BOX_BG,  C_BOX_B),
    ]
    boxes_bot = [
        (8.80, 4.20, "Evaporator B\nQ_in",      C_BOX_BG,  C_BOX_B),
        (10.80, 2.85, "Condenser\nQ_out",        C_BOX_BG,  C_BOX_B),
        (5.35,  4.20, "Pipe\n(return)",           C_BOX_BG,  C_BOX_B),
        (3.40,  4.20, "Junction /\nMixer",        C_BOX_BG,  C_BOX_B),
    ]

    def draw_box(cx, cy, label, fill, border):
        add_rect_with_text(
            slide,
            Inches(cx) - BW/2, Inches(cy) - BH/2,
            BW, BH,
            label,
            size=11, bold=False,
            fill_color=fill, line_color=border, text_color=C_DARK,
        )

    for cx, cy, lbl, f, b in boxes_top:
        draw_box(cx, cy, lbl, f, b)
    for cx, cy, lbl, f, b in boxes_bot:
        draw_box(cx, cy, lbl, f, b)

    # Splitter node (circle-ish) between supply pipe and two evaporators
    SX = Inches(7.00); SY = Inches(2.85)
    add_rect_with_text(slide, SX - Inches(0.22), SY - Inches(0.22),
                       Inches(0.44), Inches(0.44),
                       "⬤", size=11, fill_color=C_LGRAY,
                       line_color=C_GRAY, text_color=C_GRAY)

    # -- Arrows (top row, left to right) --
    arrow_y1 = Inches(1.55)
    # Acc → Pump
    add_arrow_right(slide, Inches(1.45+0.75), arrow_y1, Inches(3.40-0.75 - 1.45-0.75 - 0.05))
    # label: mdot, ΔP_p
    add_textbox(slide, Inches(2.00), arrow_y1 - Inches(0.38), Inches(1.1), Inches(0.30),
                "ṁ, ΔP_p", size=9, color=C_BLUE, align=PP_ALIGN.CENTER)

    # Pump → Pipe supply
    add_arrow_right(slide, Inches(3.40+0.75), arrow_y1, Inches(5.35-0.75 - (3.40+0.75)))
    add_textbox(slide, Inches(3.85), arrow_y1 - Inches(0.38), Inches(1.1), Inches(0.30),
                "ṁ, h, P", size=9, color=C_BLUE, align=PP_ALIGN.CENTER)

    # Pipe → Splitter node
    add_arrow_right(slide, Inches(5.35+0.75), arrow_y1, Inches(7.00 - (5.35+0.75)))
    # Splitter → Eva A
    add_arrow_right(slide, Inches(7.00+0.22), arrow_y1, Inches(8.80-0.75 - (7.00+0.22)))
    # Splitter → Eva B (vertical down, then right handled by connector)
    add_arrow_down(slide, Inches(7.00), Inches(2.85+0.22), Inches(4.20 - 2.85 - 0.22 - 0.22))
    add_arrow_right(slide, Inches(7.00+0.22), Inches(4.20),
                    Inches(8.80-0.75 - (7.00+0.22)))

    # Eva A → Condenser (vertical connector)
    add_arrow_down(slide, Inches(8.80+0.75 - 0.35), arrow_y1 + Inches(BH/2),
                   Inches(2.85 - 1.55 - BH/2))
    add_arrow_right(slide, Inches(8.80+0.75 - 0.35), Inches(2.85),
                    Inches(10.80-0.75 - (8.80+0.75-0.35) - 0.05))

    # Eva B → merge node
    add_arrow_right(slide, Inches(8.80+0.75), Inches(4.20),
                    Inches(10.80-0.75 - (8.80+0.75) - 0.05))
    # merge (use condenser box centre x for simplicity)
    add_arrow_down(slide, Inches(10.80 - 0.35), Inches(2.85+0.28),
                   Inches(4.20 - 2.85 - 0.28 - 0.28))

    # Condenser → Return pipe (go left at bottom row)
    add_arrow_right(slide, Inches(5.35+0.75), Inches(4.20),
                    Inches(3.40+0.75 - (5.35+0.75) - 0.05))  # wrong dir; use explicit
    # Return pipe ← Condenser: draw right-to-left as left connector from condenser
    cx1 = slide.shapes.add_connector(1,
        Inches(10.80-0.75), Inches(4.20),
        Inches(5.35+0.75), Inches(4.20))
    cx1.line.color.rgb = C_GRAY
    cx1.line.width = Pt(1.5)

    # Return pipe → Mixer
    cx2 = slide.shapes.add_connector(1,
        Inches(5.35-0.75), Inches(4.20),
        Inches(3.40+0.75), Inches(4.20))
    cx2.line.color.rgb = C_GRAY
    cx2.line.width = Pt(1.5)

    # Mixer → Accumulator (go up)
    cx3 = slide.shapes.add_connector(1,
        Inches(3.40-0.35), Inches(4.20),
        Inches(1.45-0.35), Inches(1.55))
    cx3.line.color.rgb = C_TEAL
    cx3.line.width = Pt(1.5)

    # Variable labels
    add_textbox(slide, Inches(9.3), Inches(0.85), Inches(1.8), Inches(0.28),
                "Q_out → sink", size=9, italic=True, color=C_GRAY)
    add_textbox(slide, Inches(9.3), Inches(1.35), Inches(1.8), Inches(0.28),
                "(condenser)", size=9, italic=True, color=C_GRAY)
    add_textbox(slide, Inches(0.4), Inches(2.60), Inches(1.5), Inches(0.28),
                "P_ref", size=10, bold=True, color=C_TEAL)

    # Key variables box (right side)
    add_rect(slide, Inches(11.6), Inches(1.1), Inches(1.5), Inches(2.8),
             fill_color=C_BOX_BG, line_color=C_BOX_B)
    add_multiline_textbox(
        slide,
        Inches(11.65), Inches(1.15),
        Inches(1.4), Inches(2.7),
        [
            ("Variables", 11, True, False, C_BLUE, PP_ALIGN.LEFT),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("P   pressure", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("h   enthalpy", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("ṁ   mass flow", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("x   quality", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Q_in  heat load", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Q_out heat rejection", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("ΔP_p  pump rise", 10, False, False, C_DARK, PP_ALIGN.LEFT),
        ],
    )

    # Caption
    add_textbox(slide, Inches(0.3), Inches(5.35), SW - Inches(0.6), Inches(0.55),
        "Loop pressure closure: the sum of pressure changes around any closed path must equal zero. "
        "Phase change occurs in the evaporator (two-phase) and condenser (condensation). "
        "Parallel branches share the same ΔP; branch flow rates are solver-determined.",
        size=11, italic=True, color=C_GRAY, align=PP_ALIGN.LEFT,
    )

    add_footer(slide, 2)
    set_notes(slide, (
        "This diagram shows the essential topology of a single-loop MPL "
        "with two parallel evaporator branches and one condenser.\n\n"
        "The key physical coupling is the pressure loop closure: the pump provides a pressure rise "
        "that must exactly balance the sum of all distributed and localised pressure drops. "
        "The accumulator sets the absolute pressure reference.\n\n"
        "In the two-phase regions — the evaporators and part of the condenser — "
        "enthalpy, quality, heat transfer and two-phase pressure drop are tightly coupled. "
        "Parallel branches add a branch maldistribution problem: "
        "the framework must determine how the total mass flow splits between branches.\n\n"
        "The variables P, h, ṁ are the primary unknowns. "
        "Temperature, quality, density and phase are derived quantities."
    ))
    return slide


# ---------------------------------------------------------------------------
# Slide 3 — State representation and property philosophy
# ---------------------------------------------------------------------------

def slide_03(prs):
    slide = blank_slide(prs)
    title_bar(slide, "State Representation and Property Philosophy",
              "Primary unknowns, derived quantities, and the (P, h) choice")

    CT = Inches(content_top(has_subtitle=True))
    LM = Inches(0.45)
    COL1 = Inches(6.20)   # left panel width
    COL2 = Inches(6.40)   # right panel

    # Left panel — SystemState
    add_rect(slide, LM, CT, COL1, Inches(3.5),
             fill_color=C_BOX_BG, line_color=C_BOX_B)
    add_multiline_textbox(
        slide,
        LM + Inches(0.15), CT + Inches(0.10),
        COL1 - Inches(0.3), Inches(3.35),
        [
            ("SystemState  (solver-owned)", 14, True, False, C_BLUE, PP_ALIGN.LEFT),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Stored primary unknowns (flat ordered vector):", 12, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  P    pressure at each port-node", 12, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  h    specific enthalpy at each port-node", 12, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  ṁ    mass flow rate at each port-node", 12, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  V_g  accumulator gas volume  (internal state)", 12, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  ω    pump shaft speed         (named seam, frozen v1)", 12, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("ṁ is NOT part of FluidState.", 12, True, False, C_AMBER, PP_ALIGN.LEFT),
            ("Ports carry NO values — connectivity only.", 12, True, False, C_AMBER, PP_ALIGN.LEFT),
        ],
    )

    # Right panel — FluidState
    RX = LM + COL1 + Inches(0.3)
    add_rect(slide, RX, CT, COL2, Inches(3.5),
             fill_color=C_TEAL_BG, line_color=C_TEAL)
    add_multiline_textbox(
        slide,
        RX + Inches(0.15), CT + Inches(0.10),
        COL2 - Inches(0.3), Inches(3.35),
        [
            ("FluidState  (P, h, identity)", 14, True, False, C_TEAL, PP_ALIGN.LEFT),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Derived on demand — never stored:", 12, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  T, T_sat, ρ, μ, k, σ, c_p, phase", 12, False, True, C_DARK, PP_ALIGN.LEFT),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Quality (vapour fraction):", 12, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  x = (h − h_f(P)) / (h_g(P) − h_f(P))", 13, True, False, C_BLUE, PP_ALIGN.LEFT),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Continuous across subcooled / two-phase /", 12, False, False, C_DARK, PP_ALIGN.LEFT),
            ("superheated — no region-switching required.", 12, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("PropertyBackend: CoolProp (default),", 12, False, False, C_DARK, PP_ALIGN.LEFT),
            ("REFPROP, tabulated, or empirical — swappable.", 12, False, False, C_DARK, PP_ALIGN.LEFT),
        ],
    )

    # Why (P,h) box
    add_rect(slide, LM, CT + Inches(3.65), SW - LM * 2, Inches(1.30),
             fill_color=RGBColor(0xF5, 0xF5, 0xF5), line_color=C_LGRAY)
    add_multiline_textbox(
        slide,
        LM + Inches(0.15), CT + Inches(3.72),
        SW - LM * 2 - Inches(0.3), Inches(1.2),
        [
            ("Why (P, h)?", 12, True, False, C_BLUE, PP_ALIGN.LEFT),
            ("  (P, h) is single-valued and continuous across all thermodynamic regions. "
             "Storing T or ρ beside (P, h) creates a second source of truth that can drift silently. "
             "In dynamics, h is the natural energy-equation variable; adding dynamics "
             "activates stored derivatives, it does not change the state representation.", 11, False, False, C_GRAY, PP_ALIGN.LEFT),
        ],
    )

    add_footer(slide, 3)
    set_notes(slide, (
        "The state representation is a foundational design decision.\n\n"
        "SystemState is the solver-owned flat vector of all primary unknowns: "
        "pressure P, specific enthalpy h, and mass flow ṁ at every port-node, "
        "plus named internal states such as the accumulator gas volume V_g and the pump shaft speed ω.\n\n"
        "FluidState is a pure value object (P, h, identity). "
        "It is ephemeral — constructed transiently when a property is needed and then discarded. "
        "All derived quantities — temperature, density, quality, viscosity — are computed on demand "
        "from the PropertyBackend. They are never stored.\n\n"
        "The key equation is the quality relation: x = (h − h_f(P)) / (h_fg(P)). "
        "It is continuous and differentiable across the saturation dome, "
        "which is essential for Newton convergence without region-switching.\n\n"
        "Mass flow rate ṁ is part of SystemState, not FluidState. "
        "FluidState is two numbers plus identity — nothing else.\n\n"
        "Ports carry no values. They are connectivity declarations only."
    ))
    return slide


# ---------------------------------------------------------------------------
# Slide 4 — Architecture as physical responsibility separation
# ---------------------------------------------------------------------------

def slide_04(prs):
    slide = blank_slide(prs)
    title_bar(slide, "Architecture as Physical Responsibility Separation",
              "A strictly acyclic layered dependency graph — dependencies flow in one direction only")

    # Draw the layered DAG as stacked horizontal bands
    CT = Inches(1.0)
    LM = Inches(0.40)
    BW = Inches(9.20)   # width of each layer band
    BH = Inches(0.54)
    GAP = Inches(0.06)

    layers = [
        ("Layer 0 — Inert data",      "Geometry  |  PropertyBackend",
         C_TEAL_BG, C_TEAL),
        ("Layer 1 — Thermodynamic state",  "FluidState (P, h, identity) → PropertyBackend",
         C_BOX_BG, C_BOX_B),
        ("Layer 2 — Interface",        "Port  (connectivity only — carries no numerical values)",
         RGBColor(0xEF, 0xF0, 0xFA), RGBColor(0x8A, 0x8A, 0xCC)),
        ("Layer 3 — Closure",          "Correlation  (pure, stateless:  evaluate(input) → value + validity)",
         RGBColor(0xFB, 0xF0, 0xE8), RGBColor(0xCC, 0x8A, 0x50)),
        ("Layer 4 — Modifier",         "Calibration  (named multipliers on closure outputs — not on balances)",
         RGBColor(0xFD, 0xF6, 0xE3), RGBColor(0xC8, 0xAA, 0x30)),
        ("Layer 5 — Local physics",    "Component  (geometry + correlations + calibration → residuals)",
         RGBColor(0xF0, 0xF7, 0xEE), RGBColor(0x60, 0xA0, 0x5A)),
        ("Layer 6 — Topology",         "Network  (topology, loop closure, branch structure, pressure-reference wiring)",
         RGBColor(0xF2, 0xEC, 0xF8), RGBColor(0x90, 0x6A, 0xB8)),
        ("Layer 7 — Numerics",         "Solver  (assembles residuals, drives convergence — contains NO physics)",
         RGBColor(0xF7, 0xEC, 0xEC), RGBColor(0xB8, 0x5A, 0x5A)),
        ("Output",                     "Result  (converged state + invariant residuals + calibration report + tuple reference)",
         RGBColor(0xE8, 0xF2, 0xE8), RGBColor(0x50, 0x90, 0x50)),
    ]

    for i, (layer_label, desc, fill, border) in enumerate(layers):
        y = CT + i * (BH + GAP)
        add_rect(slide, LM, y, BW, BH, fill_color=fill, line_color=border)
        # Layer label (left, bold)
        add_textbox(slide, LM + Inches(0.08), y + Inches(0.06),
                    Inches(2.30), BH - Inches(0.1),
                    layer_label, size=10, bold=True, color=C_DARK)
        # Description (right)
        add_textbox(slide, LM + Inches(2.40), y + Inches(0.06),
                    BW - Inches(2.50), BH - Inches(0.1),
                    desc, size=10, color=C_GRAY)
        # Arrow downward between layers (except last)
        if i < len(layers) - 1:
            ax = LM + BW + Inches(0.10)
            ay = y + BH / 2
            # no arrows drawn on right — the stacking itself implies direction

    # Arrow label on the right
    arrow_top    = CT
    arrow_bottom = CT + len(layers) * (BH + GAP) - GAP
    arr_x = LM + BW + Inches(0.15)
    cx = slide.shapes.add_connector(1, arr_x, arrow_top, arr_x, arrow_bottom)
    cx.line.color.rgb = C_BLUE
    cx.line.width = Pt(2.0)
    add_textbox(slide, arr_x + Inches(0.05), arrow_top + Inches(0.5),
                Inches(1.5), Inches(1.0),
                "dependencies\nflow down\nonly", size=9, color=C_BLUE,
                align=PP_ALIGN.LEFT)

    # Right side note
    note_x = LM + BW + Inches(0.15)
    note_y = CT + Inches(3.5)
    add_rect(slide, note_x, note_y, Inches(2.80), Inches(2.40),
             fill_color=C_BOX_BG, line_color=C_BOX_B)
    add_multiline_textbox(
        slide,
        note_x + Inches(0.1), note_y + Inches(0.08),
        Inches(2.6), Inches(2.25),
        [
            ("Three sanctioned seams:", 11, True, False, C_BLUE, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("① Swap a Correlation\n   → config change only", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("② Edit the Network topology\n   → topology change only", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("③ Swap the Solver\n   → numerics change only", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("The DAG guarantees none\ndisturbs the others.", 10, False, True, C_GRAY, PP_ALIGN.LEFT),
        ],
    )

    add_footer(slide, 4)
    set_notes(slide, (
        "The architecture is a strictly acyclic layered dependency graph. "
        "Every concept has a layer, and dependencies are permitted only downward.\n\n"
        "Layer 0 is inert data: Geometry and PropertyBackend. "
        "They know nothing above them.\n\n"
        "Layer 1 is FluidState. It reads the PropertyBackend and nothing else.\n\n"
        "Ports at Layer 2 are connectivity declarations only. "
        "They carry no numerical values.\n\n"
        "Correlations at Layer 3 are pure, stateless closure relations. "
        "They receive a typed input struct and return a value plus a validity verdict. "
        "They do not know the component, the network, or the solver.\n\n"
        "Calibration at Layer 4 applies named multipliers to closure outputs — "
        "never to the conservation equations themselves.\n\n"
        "Components at Layer 5 own local physics. "
        "They contribute residuals and derivatives. "
        "They do not know their neighbours, the network, or the solver.\n\n"
        "The Network at Layer 6 states what must hold: "
        "topology, loop closure, branch structure, pressure-reference wiring. "
        "The Solver at Layer 7 determines how to satisfy it numerically.\n\n"
        "This layering is not software elegance for its own sake. "
        "It makes every physical assumption explicit, replaceable, and independently validatable."
    ))
    return slide


# ---------------------------------------------------------------------------
# Slide 5 — Governing balances and closure seams
# ---------------------------------------------------------------------------

def slide_05(prs):
    slide = blank_slide(prs)
    title_bar(slide, "Governing Balances and Closure Seams",
              "Conservation equations vs. closure relations — a strict distinction")

    CT = Inches(content_top(has_subtitle=True))
    LM = Inches(0.45)
    HALF = (SW - LM * 2 - Inches(0.3)) / 2

    # Left: conservation equations
    add_rect(slide, LM, CT, HALF + Inches(0.1), Inches(4.30),
             fill_color=C_BOX_BG, line_color=C_BOX_B)
    add_multiline_textbox(
        slide,
        LM + Inches(0.15), CT + Inches(0.10),
        HALF - Inches(0.05), Inches(4.20),
        [
            ("Conservation equations  (never scaled)", 13, True, False, C_BLUE, PP_ALIGN.LEFT),
            ("", 6, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Mass balance at each junction:", 11, True, False, C_DARK, PP_ALIGN.LEFT),
            ("  Σ ṁ_in  −  Σ ṁ_out  =  0", 13, True, False, C_DARK, PP_ALIGN.LEFT),
            ("", 6, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Energy balance (steady, per component):", 11, True, False, C_DARK, PP_ALIGN.LEFT),
            ("  Q̇ − Ẇ + Σ ṁ_in h_in − Σ ṁ_out h_out = 0", 13, True, False, C_DARK, PP_ALIGN.LEFT),
            ("", 6, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Pressure loop closure (per closed path):", 11, True, False, C_DARK, PP_ALIGN.LEFT),
            ("  Σ ΔP  =  0", 13, True, False, C_DARK, PP_ALIGN.LEFT),
            ("", 6, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Pipe pressure gradient (per cell, 1D):", 11, True, False, C_DARK, PP_ALIGN.LEFT),
            ("  dP/dx  =  (dP/dx)_fric  +  (dP/dx)_grav  +  (dP/dx)_acc", 12, True, False, C_DARK, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("   = R* · (dP/dx)_fric,corr  +  ρg(dz/dx)  +  d(G²v)/dx", 11, False, True, C_GRAY, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("R* = calibration multiplier (friction term only)", 10, False, True, C_AMBER, PP_ALIGN.LEFT),
        ],
    )

    # Right: closure quantities
    RX = LM + HALF + Inches(0.3)
    add_rect(slide, RX, CT, HALF + Inches(0.1), Inches(4.30),
             fill_color=C_TEAL_BG, line_color=C_TEAL)
    add_multiline_textbox(
        slide,
        RX + Inches(0.15), CT + Inches(0.10),
        HALF - Inches(0.05), Inches(4.20),
        [
            ("Closure quantities  (replaceable, validated)", 13, True, False, C_TEAL, PP_ALIGN.LEFT),
            ("", 6, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Single-phase friction gradient:", 11, True, False, C_DARK, PP_ALIGN.LEFT),
            ("  (dP/dx)_fric  from Churchill (1977),", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  Müller-Steinhagen-Heck, or user-supplied", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 6, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Two-phase pressure drop:", 11, True, False, C_DARK, PP_ALIGN.LEFT),
            ("  Friedel, Müller-Steinhagen-Heck, …", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  (Phase 11 — in progress)", 11, False, True, C_GRAY, PP_ALIGN.LEFT),
            ("", 6, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Heat transfer coefficient:", 11, True, False, C_DARK, PP_ALIGN.LEFT),
            ("  Shah, Gungor-Winterton, Kim-Mudawar, …", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  (Phase 11 — in progress)", 11, False, True, C_GRAY, PP_ALIGN.LEFT),
            ("", 6, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Void fraction / two-phase density:", 11, True, False, C_DARK, PP_ALIGN.LEFT),
            ("  HEM-type admissible; framework is", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  closure-agnostic — slip, drift-flux,", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  regime-dependent closures are future options", 11, False, False, C_DARK, PP_ALIGN.LEFT),
        ],
    )

    # Bottom bar: calibration rule
    add_rect(slide, LM, CT + Inches(4.45), SW - LM * 2, Inches(0.78),
             fill_color=RGBColor(0xFD, 0xF8, 0xEC), line_color=C_AMBER)
    add_multiline_textbox(
        slide,
        LM + Inches(0.15), CT + Inches(4.50),
        SW - LM * 2 - Inches(0.3), Inches(0.70),
        [
            ("Calibration firewall:", 12, True, False, C_AMBER, PP_ALIGN.LEFT),
            ("  Calibration multipliers scale closure outputs only (friction gradient, HTC). "
             "They never scale conservation equations. "
             "A wrong calibration worsens the balance residual — it cannot produce a false-passing balance.", 11, False, False, C_DARK, PP_ALIGN.LEFT),
        ],
    )

    add_footer(slide, 5)
    set_notes(slide, (
        "This slide makes the most important physical distinction in the framework explicit:\n\n"
        "Conservation equations — mass, energy, and pressure loop closure — "
        "are inviolable physical laws. They are assembled by the solver "
        "from component contributions and network conditions. They are never scaled.\n\n"
        "Closure relations — friction gradient, heat transfer coefficient, void fraction — "
        "are semi-empirical substitutes for unresolved microscale physics. "
        "They are swappable, validity-bounded, and the primary research seam.\n\n"
        "The pressure gradient in a 1D passage is decomposed into three terms: "
        "friction, gravity, and acceleration. "
        "The calibration multiplier R* applies only to the friction term. "
        "Gravity and acceleration are first-principles physics — they are never fudged.\n\n"
        "On void fraction and two-phase closures: "
        "HEM-type closures are admissible and may be used in early two-phase models, "
        "but the framework is closure-agnostic. "
        "Slip-flow, drift-flux, regime-dependent, and empirical void-fraction closures "
        "can be substituted through the same correlation slot without changing the conservation equations.\n\n"
        "Single-phase friction is implemented via Churchill (1977). "
        "Two-phase DP and HTC correlations are in progress in Phase 11."
    ))
    return slide


# ---------------------------------------------------------------------------
# Slide 6 — Pump and accumulator
# ---------------------------------------------------------------------------

def slide_06(prs):
    slide = blank_slide(prs)
    title_bar(slide, "Pump and Accumulator: Drive and Pressure Reference",
              "Phase 10 — implemented at V1 fidelity")

    CT = Inches(content_top(has_subtitle=True))
    LM = Inches(0.45)
    HALF = (SW - LM * 2 - Inches(0.3)) / 2

    # Pump panel
    add_rect(slide, LM, CT, HALF + Inches(0.1), Inches(3.70),
             fill_color=C_BOX_BG, line_color=C_BOX_B)
    add_multiline_textbox(
        slide,
        LM + Inches(0.15), CT + Inches(0.10),
        HALF - Inches(0.05), Inches(3.55),
        [
            ("Pump  — pressure-rise / map / command seam", 13, True, False, C_BLUE, PP_ALIGN.LEFT),
            ("", 6, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Prescribed pressure rise:", 11, True, False, C_DARK, PP_ALIGN.LEFT),
            ("  ΔP_p  =  ΔP_setpoint · R_p", 13, True, False, C_DARK, PP_ALIGN.LEFT),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Performance-map evaluation:", 11, True, False, C_DARK, PP_ALIGN.LEFT),
            ("  ΔP_p  =  f(ṁ, ω)  from interpolated map", 12, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Power / efficiency seam:", 11, True, False, C_DARK, PP_ALIGN.LEFT),
            ("  Ẇ_p  ≈  ṁ · ΔP_p · v / η_p", 13, True, False, C_DARK, PP_ALIGN.LEFT),
            ("  v = specific volume (from FluidState)", 10, False, True, C_GRAY, PP_ALIGN.LEFT),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Shaft-speed ω: named internal state (frozen v1)", 10, False, True, C_AMBER, PP_ALIGN.LEFT),
            ("Commands: speed setpoint or flow target", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Pump does NOT know Network or Solver", 10, False, True, C_AMBER, PP_ALIGN.LEFT),
        ],
    )

    # Accumulator panel
    RX = LM + HALF + Inches(0.3)
    add_rect(slide, RX, CT, HALF + Inches(0.1), Inches(3.70),
             fill_color=C_TEAL_BG, line_color=C_TEAL)
    add_multiline_textbox(
        slide,
        RX + Inches(0.15), CT + Inches(0.10),
        HALF - Inches(0.05), Inches(3.55),
        [
            ("Accumulator  — pressure-reference component", 13, True, False, C_TEAL, PP_ALIGN.LEFT),
            ("", 6, False, False, C_DARK, PP_ALIGN.LEFT),
            ("PCA (pneumatic) pressure law:", 11, True, False, C_DARK, PP_ALIGN.LEFT),
            ("  P(V_g)  =  P_charge · (V_charge / V_g)^n", 13, True, False, C_DARK, PP_ALIGN.LEFT),
            ("  V_g : gas volume (internal state in SystemState)", 10, False, True, C_GRAY, PP_ALIGN.LEFT),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Law slot — swappable:", 11, True, False, C_DARK, PP_ALIGN.LEFT),
            ("  PCA (implemented)  |  HCA (declared seam)", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  bellows / spring / gas-charged (future)", 11, False, True, C_GRAY, PP_ALIGN.LEFT),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Ownership split:", 11, True, False, C_DARK, PP_ALIGN.LEFT),
            ("  Accumulator owns: law, value, geometry", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  Network owns: which node is the reference", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  Solver owns: global consistency", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("A second accumulator → topology validation error", 10, False, True, C_AMBER, PP_ALIGN.LEFT),
        ],
    )

    # Bottom note
    add_rect(slide, LM, CT + Inches(3.85), SW - LM * 2, Inches(0.72),
             fill_color=RGBColor(0xF8, 0xF8, 0xF8), line_color=C_LGRAY)
    add_multiline_textbox(
        slide,
        LM + Inches(0.15), CT + Inches(3.90),
        SW - LM * 2 - Inches(0.3), Inches(0.65),
        [
            ("Phase 10 status:", 12, True, False, C_GREEN, PP_ALIGN.LEFT),
            ("  Pump map/command/power seam, accumulator PCA law, V_g internal state, and Network pressure-reference wiring are implemented and audited. "
             "Dynamic shaft-speed derivatives and full loop pressure convergence are Phase 11+ work.",
             11, False, False, C_GRAY, PP_ALIGN.LEFT),
        ],
    )

    add_footer(slide, 6)
    set_notes(slide, (
        "Phase 10 is complete. These are the two components that drive and anchor the MPL loop.\n\n"
        "The pump is modelled as a pressure-rise component. "
        "In V1 it supports a prescribed pressure-rise setpoint and interpolation from a performance map. "
        "The power/efficiency seam computes shaft power from mass flow, pressure rise, "
        "specific volume, and mechanical efficiency. "
        "Shaft speed ω is a named internal state; in V1 it is frozen (zero derivative). "
        "Dynamics will unfreeze it — no new state variable will be needed.\n\n"
        "The accumulator is a pressure-reference component. "
        "For a pneumatic compensation accumulator (PCA), the pressure law is a polytropic relation "
        "between gas volume V_g and pressure. "
        "V_g is the stored internal state; pressure is derived from it.\n\n"
        "The pressure-reference ownership is split deliberately: "
        "the accumulator owns the law and the value; "
        "the network owns which node is the reference and enforces that there is exactly one; "
        "the solver owns global consistency. "
        "This means a second accumulator is caught at topology validation, not as a numerical pathology.\n\n"
        "HCA and other law types are declared as future seams. "
        "PCA is the implemented V1 closure."
    ))
    return slide


# ---------------------------------------------------------------------------
# Slide 7 — HeatExchangerModel, Evaporator, Condenser
# ---------------------------------------------------------------------------

def slide_07(prs):
    slide = blank_slide(prs)
    title_bar(slide, "Next Development: HeatExchangerModel, Evaporator and Condenser",
              "Phase 11 — in progress (foundation and ε-NTU checkpoints merged)")

    CT = Inches(content_top(has_subtitle=True))
    LM = Inches(0.45)

    # Central HXModel diagram
    # Inputs on the left, box in the middle, outputs on the right
    BOX_X = Inches(4.40)
    BOX_Y = CT + Inches(0.20)
    BOX_W = Inches(3.20)
    BOX_H = Inches(2.40)

    add_rect(slide, BOX_X, BOX_Y, BOX_W, BOX_H,
             fill_color=C_BOX_BG, line_color=C_BLUE)
    add_multiline_textbox(
        slide,
        BOX_X + Inches(0.12), BOX_Y + Inches(0.10),
        BOX_W - Inches(0.24), BOX_H - Inches(0.20),
        [
            ("HeatExchangerModel", 13, True, False, C_BLUE, PP_ALIGN.CENTER),
            ("(component-internal strategy)", 10, False, True, C_GRAY, PP_ALIGN.CENTER),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Strategies (not correlations):", 11, True, False, C_DARK, PP_ALIGN.LEFT),
            ("  ε-NTU  (implemented — foundation+11B)", 10, False, False, C_GREEN, PP_ALIGN.LEFT),
            ("  LMTD  (planned)", 10, False, True, C_GRAY, PP_ALIGN.LEFT),
            ("  Segmented marching  (planned)", 10, False, True, C_GRAY, PP_ALIGN.LEFT),
            ("  Moving boundary  (declared seam)", 10, False, True, C_GRAY, PP_ALIGN.LEFT),
        ],
    )

    # Arrows left → box
    arr_y_mid = BOX_Y + BOX_H / 2
    add_arrow_right(slide, Inches(0.45), arr_y_mid, BOX_X - Inches(0.45) - Inches(0.1))

    # Input labels
    input_lines = [
        "PropertyBackend  (fluid properties)",
        "HTC correlations  (Shah, KM, …)",
        "ΔP correlations  (Friedel, …)",
        "Geometry  (plate, microchannel, …)",
        "Secondary BC  (T_sink, ṁ_sink, cp_sink)",
        "Discretization  (Lumped / Segmented)",
    ]
    for k, lbl in enumerate(input_lines):
        add_textbox(slide, LM, CT + Inches(0.25) + k * Inches(0.40),
                    Inches(3.80), Inches(0.38),
                    "→  " + lbl, size=10, color=C_DARK)

    # Arrows box → right
    OUT_X = BOX_X + BOX_W + Inches(0.12)
    add_arrow_right(slide, BOX_X + BOX_W, arr_y_mid,
                    SW - (BOX_X + BOX_W) - Inches(0.50))

    output_lines = [
        "Q̇  (heat rate)",
        "h_out  (outlet enthalpy)",
        "ΔP  (pressure drop)",
        "residuals",
        "validity verdicts",
        "calibration report",
    ]
    for k, lbl in enumerate(output_lines):
        add_textbox(slide, OUT_X + Inches(0.18), CT + Inches(0.25) + k * Inches(0.40),
                    Inches(3.0), Inches(0.38),
                    lbl, size=10, color=C_DARK)

    # Bottom section: key distinctions
    SEP_Y = CT + BOX_H + Inches(0.50)
    add_rect(slide, LM, SEP_Y, SW - LM * 2, Inches(0.04),
             fill_color=C_LGRAY)

    add_multiline_textbox(
        slide,
        LM, SEP_Y + Inches(0.12),
        SW - LM * 2, Inches(1.0),
        [
            ("Key distinction:", 12, True, False, C_BLUE, PP_ALIGN.LEFT),
            ("  A Correlation returns one local closure value. "
             "A HeatExchangerModel is a solution strategy that consumes correlations, geometry, "
             "secondary boundary conditions and the PropertyBackend to produce an integrated thermal result. "
             "ε-NTU is a HeatExchangerModel strategy — not a correlation.",
             11, False, False, C_DARK, PP_ALIGN.LEFT),
        ],
    )

    # Parallel branch note
    add_rect(slide, LM, SEP_Y + Inches(1.20), SW - LM * 2, Inches(0.90),
             fill_color=RGBColor(0xF4, 0xF8, 0xF4), line_color=C_LGRAY)
    add_multiline_textbox(
        slide,
        LM + Inches(0.15), SEP_Y + Inches(1.26),
        SW - LM * 2 - Inches(0.3), Inches(0.80),
        [
            ("Scientific opportunity — parallel evaporators and branch flow distribution:", 12, True, False, C_GREEN, PP_ALIGN.LEFT),
            ("  Each branch carries its own HeatExchangerModel, geometry, HTC/DP closures, and flow condition. "
             "Branch flows are solver-determined from equal-ΔP conditions. "
             "This framework supports maldistribution studies natively through topology — no special-casing.",
             11, False, False, C_DARK, PP_ALIGN.LEFT),
        ],
    )

    add_footer(slide, 7)
    set_notes(slide, (
        "Phase 11 introduces the HeatExchangerModel, Evaporator, and Condenser.\n\n"
        "The critical architectural distinction is that a HeatExchangerModel is not a Correlation. "
        "A Correlation returns a single local closure value — a friction factor, an HTC, a void fraction. "
        "A HeatExchangerModel is a complete solution strategy for a heat exchanger: "
        "it consumes HTC correlations, DP correlations, geometry, the PropertyBackend, "
        "and secondary boundary conditions, and it produces an integrated result — "
        "heat rate, outlet enthalpy, pressure drop, and residuals.\n\n"
        "ε-NTU is implemented as a HeatExchangerModel strategy. "
        "It is not a correlation and does not belong in the correlation registry.\n\n"
        "What is implemented so far: the HeatExchangerModel contract and registry, "
        "EvaporatorComponent and CondenserComponent wrappers, "
        "and the ε-NTU model with fixed-heat-rate and sink-inlet-temperature-and-flow modes "
        "including explicit PrimaryThermalMode and UAComputationMode.\n\n"
        "LMTD, segmented marching, and moving-boundary models are planned but not yet implemented.\n\n"
        "The parallel evaporator case is particularly important scientifically: "
        "the framework handles it through topology — adding a branch is a network edit, "
        "not a solver or component change. This makes maldistribution studies straightforward."
    ))
    return slide


# ---------------------------------------------------------------------------
# Slide 8 — Scientific opportunities, limitations and validation roadmap
# ---------------------------------------------------------------------------

def slide_08(prs):
    slide = blank_slide(prs)
    title_bar(slide, "Scientific Opportunities, Limitations and Validation Roadmap",
              "\"Architecture is used here to make physical assumptions explicit, replaceable and testable.\"")

    CT = Inches(content_top(has_subtitle=True))
    LM = Inches(0.45)
    THIRD = (SW - LM * 2 - Inches(0.4)) / 3

    # Column 1: Scientific opportunities
    add_rect(slide, LM, CT, THIRD, Inches(3.80),
             fill_color=RGBColor(0xF0, 0xF7, 0xEE), line_color=C_GREEN)
    add_multiline_textbox(
        slide,
        LM + Inches(0.12), CT + Inches(0.08),
        THIRD - Inches(0.20), Inches(3.70),
        [
            ("Scientific opportunities", 13, True, False, C_GREEN, PP_ALIGN.LEFT),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("•  Parallel evaporators and branch\n   maldistribution", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("•  Correlation comparison studies\n   (Shah vs Kim-Mudawar, etc.)", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("•  HEM vs slip / drift-flux /\n   void-fraction closure alternatives", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("•  Accumulator influence on loop\n   pressure stability", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("•  DOE / parametric sweeps over\n   geometry and operating point", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("•  Dynamic modelling and transient\n   response (future seams prepared)", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("•  Control-oriented linearisation\n   (MPC / ROM extraction)", 11, False, False, C_DARK, PP_ALIGN.LEFT),
        ],
    )

    # Column 2: Limitations
    COL2X = LM + THIRD + Inches(0.2)
    add_rect(slide, COL2X, CT, THIRD, Inches(3.80),
             fill_color=RGBColor(0xFD, 0xF8, 0xEC), line_color=C_AMBER)
    add_multiline_textbox(
        slide,
        COL2X + Inches(0.12), CT + Inches(0.08),
        THIRD - Inches(0.20), Inches(3.70),
        [
            ("Current limitations", 13, True, False, C_AMBER, PP_ALIGN.LEFT),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("•  Steady-state only\n   (dynamics: seams declared)", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("•  1D / lumped component contributions\n   (segmented fidelity available)", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("•  Correlations are semi-empirical\n   and validity-bounded", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("•  Two-phase HX physics (Phase 11)\n   not yet complete", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("•  HCA accumulator: declared seam,\n   not yet implemented", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("•  LMTD, segmented-march, moving-\n   boundary HX: planned, not built", 11, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 4, False, False, C_DARK, PP_ALIGN.LEFT),
            ("•  Moving-boundary condenser:\n   dynamic event detection deferred", 11, False, False, C_DARK, PP_ALIGN.LEFT),
        ],
    )

    # Column 3: Validation roadmap
    COL3X = COL2X + THIRD + Inches(0.2)
    add_rect(slide, COL3X, CT, THIRD, Inches(3.80),
             fill_color=C_BOX_BG, line_color=C_BOX_B)
    add_multiline_textbox(
        slide,
        COL3X + Inches(0.12), CT + Inches(0.08),
        THIRD - Inches(0.20), Inches(3.70),
        [
            ("Validation roadmap", 13, True, False, C_BLUE, PP_ALIGN.LEFT),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Completed:", 11, True, False, C_GREEN, PP_ALIGN.LEFT),
            ("  Phases 1–10 audited\n  2 301 tests passing", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  Unit: FluidState, Correlation,\n  Calibration, Geometry, Pump, PCA", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  Network topology and assembly", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  Steady solver residuals", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("In progress:", 11, True, False, C_AMBER, PP_ALIGN.LEFT),
            ("  HX model + ε-NTU (Phase 11)", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("", 5, False, False, C_DARK, PP_ALIGN.LEFT),
            ("Planned:", 11, True, False, C_BLUE, PP_ALIGN.LEFT),
            ("  Kokate (2024) R-134a loop data", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  Li et al. (2021) Acetone data", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  Fujii (2004) validation data", 10, False, False, C_DARK, PP_ALIGN.LEFT),
            ("  MPL2030 experimental data\n  (future phases)", 10, False, False, C_DARK, PP_ALIGN.LEFT),
        ],
    )

    # Closing statement
    add_rect(slide, LM, CT + Inches(3.95), SW - LM * 2, Inches(0.75),
             fill_color=C_BLUE, line_color=None)
    add_textbox(
        slide,
        LM + Inches(0.20), CT + Inches(4.03),
        SW - LM * 2 - Inches(0.4), Inches(0.65),
        "\"Architecture is used here to make physical assumptions explicit, replaceable and testable.\"",
        size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
    )

    add_footer(slide, 8)
    set_notes(slide, (
        "Let me close with an honest assessment of where we are and where we are going.\n\n"
        "Scientific opportunities: "
        "The framework's explicit seams make it well-suited for correlation comparison studies — "
        "switching from Shah to Kim-Mudawar in an evaporator is a configuration change. "
        "Parallel evaporator branches with maldistribution are handled through topology. "
        "Accumulator stability, HEM versus slip-flow alternatives, "
        "and future dynamic and control-oriented work are all natural extensions.\n\n"
        "Limitations — we are honest about these:\n"
        "We are steady-state only today. Dynamic seams are prepared but not activated.\n"
        "Two-phase heat exchanger physics — boiling HTC, condensation HTC, two-phase DP — "
        "are in progress in Phase 11 and are not yet validated against literature data.\n"
        "HCA is a declared seam, not an implemented model.\n"
        "LMTD and segmented-march strategies are planned but not built.\n\n"
        "Validation roadmap: "
        "Phases 1 through 10 are audited and 2301 tests pass. "
        "Literature validation against Kokate (2024) R-134a data, Li et al. (2021), "
        "and Fujii (2004) is planned for Phase 12 and beyond. "
        "MPL2030 experimental data will enter as it becomes available.\n\n"
        "The closing statement says it directly: "
        "the architecture is not a software goal in itself. "
        "It is the mechanism by which physical assumptions become explicit, testable, and replaceable."
    ))
    return slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = new_prs()

    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")
    for i, s in enumerate(prs.slides, 1):
        # Extract title from first large textbox
        texts = [sh.text_frame.paragraphs[0].text
                 for sh in s.shapes
                 if sh.has_text_frame and sh.text_frame.paragraphs]
        title_text = next((t for t in texts if t.strip()), "(no title)")
        print(f"  Slide {i}: {title_text[:80]}")


if __name__ == "__main__":
    main()
